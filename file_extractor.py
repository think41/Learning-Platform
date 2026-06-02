import os
from typing import Tuple

MAX_CHARS = 15000


def extract_text(file_path: str) -> Tuple[str, str]:
    if not os.path.exists(file_path):
        raise FileNotFoundError(f"File not found: {file_path}")

    filename = os.path.basename(file_path)
    ext = os.path.splitext(file_path)[1].lower()

    extractors = {
        '.pdf': _from_pdf,
        '.pptx': _from_pptx,
        '.ppt': _from_pptx,
        '.docx': _from_docx,
        '.doc': _from_docx,
        '.txt': _from_txt,
        '.md': _from_txt,
    }

    extractor = extractors.get(ext)
    if not extractor:
        raise ValueError(f"Unsupported file type: {ext}. Supported: {', '.join(extractors)}")

    text = extractor(file_path)

    truncated = False
    if len(text) > MAX_CHARS:
        text = text[:MAX_CHARS]
        truncated = True

    return text, filename, truncated


def _from_pdf(path: str) -> str:
    from pypdf import PdfReader
    reader = PdfReader(path)
    return "\n".join(page.extract_text() or "" for page in reader.pages).strip()


def _from_pptx(path: str) -> str:
    from pptx import Presentation
    prs = Presentation(path)
    lines = []
    for i, slide in enumerate(prs.slides, 1):
        lines.append(f"\n--- Slide {i} ---")
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                lines.append(shape.text.strip())
    return "\n".join(lines).strip()


def _from_docx(path: str) -> str:
    from docx import Document
    doc = Document(path)
    return "\n".join(p.text for p in doc.paragraphs if p.text.strip()).strip()


def _from_txt(path: str) -> str:
    with open(path, 'r', encoding='utf-8') as f:
        return f.read()
