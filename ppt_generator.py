import os
from typing import Dict, Any
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN


THEME = {
    "bg": RGBColor(0x1E, 0x1E, 0x2E),
    "title": RGBColor(0xCD, 0xD6, 0xF4),
    "body": RGBColor(0xBA, 0xC2, 0xDE),
    "accent": RGBColor(0x89, 0xB4, 0xFA),
}


def _set_bg(slide, color: RGBColor):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _styled_textbox(slide, left, top, width, height, text, font_size, bold=False, color=None, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color or THEME["body"]
    return txBox


def generate_ppt(ppt_data: Dict[Any, Any], output_dir: str = "output") -> str:
    os.makedirs(output_dir, exist_ok=True)

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    module_title = ppt_data.get("module_title", "Module")
    slides_data = ppt_data.get("slides", [])

    # Title slide
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    _set_bg(slide, THEME["bg"])
    _styled_textbox(slide, Inches(1), Inches(2.5), Inches(11), Inches(1.2),
                    module_title, 40, bold=True, color=THEME["title"], align=PP_ALIGN.CENTER)
    _styled_textbox(slide, Inches(1), Inches(3.9), Inches(11), Inches(0.6),
                    "AI Learning Platform", 18, color=THEME["accent"], align=PP_ALIGN.CENTER)

    # Content slides
    for slide_data in slides_data:
        slide = prs.slides.add_slide(blank)
        _set_bg(slide, THEME["bg"])

        # Slide title
        _styled_textbox(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.8),
                        slide_data.get("title", ""), 28, bold=True, color=THEME["title"])

        # Divider line (thin rectangle)
        line = slide.shapes.add_shape(1, Inches(0.5), Inches(1.2), Inches(12), Pt(2))
        line.fill.solid()
        line.fill.fore_color.rgb = THEME["accent"]
        line.line.fill.background()

        # Bullet points
        content = slide_data.get("content", [])
        bullet_text = "\n".join(f"• {point}" for point in content)
        _styled_textbox(slide, Inches(0.5), Inches(1.4), Inches(12), Inches(5.5),
                        bullet_text, 18, color=THEME["body"])

        # Speaker notes
        notes = slide_data.get("speaker_notes", "")
        if notes:
            slide.notes_slide.notes_text_frame.text = notes

    # Save
    safe_title = "".join(c if c.isalnum() or c in " _-" else "_" for c in module_title)
    filename = f"{safe_title.replace(' ', '_').lower()}.pptx"
    filepath = os.path.join(output_dir, filename)
    prs.save(filepath)
    return filepath
